from docx import Document
from pptx import Presentation

doc_path = "C:/Users/Sathvik/Downloads/23eg110c33.docx"
ppt_path = "generated_presentation.pptx"

doc = Document(doc_path)
prs = Presentation()

title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

title_slide = prs.slides.add_slide(title_slide_layout)
title_slide.shapes.title.text = doc.paragraphs[0].text
title_slide.placeholders[1].text = doc.paragraphs[2].text

slide = None
content = []

for para in doc.paragraphs[3:]:
    text = para.text.strip()

    if not text:
        continue

    if para.style.name.startswith("Heading"):
        if slide and content:
            tf = slide.placeholders[1].text_frame
            tf.text = content[0]
            for line in content[1:]:
                tf.add_paragraph().text = line
        slide = prs.slides.add_slide(content_slide_layout)
        slide.shapes.title.text = text
        content = []
    else:
        content.append(text)

if slide and content:
    tf = slide.placeholders[1].text_frame
    tf.text = content[0]
    for line in content[1:]:
        tf.add_paragraph().text = line

prs.save(ppt_path)
print("PPT generated:", ppt_path)